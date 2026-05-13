"""
One-shot script : uniformise la taille des chart frames Q1-Q10 dans
templates/analysis-block.pptx à 22.58 × 11.17 cm (format Q4 actuel).

POURQUOI :
    Feedback Inès (13/05) — "Peut-on mettre tous les camemberts en 11 x 22.
    Ainsi ils sont tous identiques, qu'il y ait des commentaires ou pas."
    Q4 (slide index 6 dans le template = slide 9 dans l'output final) avait
    déjà la bonne taille (22.58 × 11.17). Les 9 autres (Q1-Q3, Q5-Q10)
    étaient à 22.58 × 15.05 — donc plus hauts, laissant moins d'espace pour
    les commentaires en dessous et créant l'effet visuel non-uniforme.

USAGE :
    python scripts/patch_analysis_block_chart_size.py

    Idempotent : si la taille est déjà à 11.17 cm, on skip.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "templates" / "analysis-block.pptx"

# Cible : format Q4 actuel (mesuré 13/05)
TARGET_WIDTH_CM = 22.58
TARGET_HEIGHT_CM = 11.17

# 1 cm = 360000 EMU
CM_TO_EMU = 360000
TARGET_WIDTH_EMU = int(TARGET_WIDTH_CM * CM_TO_EMU)
TARGET_HEIGHT_EMU = int(TARGET_HEIGHT_CM * CM_TO_EMU)

# Tolérance pour éviter les re-patches micro-imprécis
TOLERANCE_EMU = 10000  # 0.028 cm

# Q1-Q10 = slides 3 à 12 dans analysis-block.pptx (0-based)
Q_SLIDE_INDICES = range(3, 13)


def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template introuvable : {TEMPLATE}")

    prs = Presentation(str(TEMPLATE))
    changed = 0
    skipped = 0

    for idx in Q_SLIDE_INDICES:
        slide = prs.slides[idx]
        q_num = idx - 2  # Q1 = slide 3, Q2 = slide 4, ...
        chart_shape = None
        for sh in slide.shapes:
            if sh.has_chart:
                chart_shape = sh
                break
        if chart_shape is None:
            print(f"  [WARN] S{idx} (Q{q_num}) : aucun chart trouvé, skip")
            continue

        w_before = chart_shape.width
        h_before = chart_shape.height
        already_ok = (
            abs(w_before - TARGET_WIDTH_EMU) <= TOLERANCE_EMU
            and abs(h_before - TARGET_HEIGHT_EMU) <= TOLERANCE_EMU
        )
        if already_ok:
            print(f"  [SKIP] S{idx} (Q{q_num}) : déjà à {w_before / CM_TO_EMU:.2f} × {h_before / CM_TO_EMU:.2f} cm")
            skipped += 1
            continue

        chart_shape.width = Emu(TARGET_WIDTH_EMU)
        chart_shape.height = Emu(TARGET_HEIGHT_EMU)
        print(
            f"  [OK]   S{idx} (Q{q_num}) : {w_before / CM_TO_EMU:.2f} × {h_before / CM_TO_EMU:.2f} cm "
            f"→ {TARGET_WIDTH_CM:.2f} × {TARGET_HEIGHT_CM:.2f} cm"
        )
        changed += 1

    if changed > 0:
        prs.save(str(TEMPLATE))
        print(f"\nTemplate sauvé : {TEMPLATE}")
        print(f"Slides modifiées : {changed}  |  Slides déjà OK : {skipped}")
    else:
        print(f"\nRien à faire. Slides déjà OK : {skipped}")


if __name__ == "__main__":
    main()
