"""Idempotent : supprime la 2ᵉ série fantôme "Colonne1" des 10 charts pie3D
du template `templates/analysis-block.pptx`.

Contexte : le template a été dérivé du PPTX Camarines original qui contenait,
en plus de la série TOTAL (counts par catégorie), une série "Colonne1" (proportions)
servant à un calcul Excel intermédiaire. Cette série n'est jamais mise à jour par
le générateur (qui n'update que le 1er <c:ser>) → leftover Camarines visible dans
les livrables (anneau parasite, étiquettes incohérentes).

Cf. PIC-FORMATION/PPTX-TEMPLATE-SPEC.md et postmortem audit Inès 2026-05-29.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from lxml import etree

NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
NS = {"c": NS_C}

TEMPLATE = Path(__file__).resolve().parent.parent / "templates" / "analysis-block.pptx"


def drop_extra_pie3d_series(template_path: Path) -> int:
    """Pour chaque pie3DChart, garde seulement la 1ʳᵉ série, supprime les autres.

    Retourne le nombre de séries supprimées au total.
    """
    prs = Presentation(template_path)
    dropped = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            cs = shape.chart._chartSpace
            for pie in cs.findall(".//c:pie3DChart", NS):
                sers = pie.findall("c:ser", NS)
                if len(sers) <= 1:
                    continue
                # Garder uniquement la 1ʳᵉ (idx=0, série "TOTAL")
                for extra in sers[1:]:
                    tx_vals = [e.text for e in extra.findall(".//c:tx//c:v", NS)]
                    pie.remove(extra)
                    dropped += 1
                    print(
                        f"  slide {slide_idx} | {shape.name} | drop ser {tx_vals}"
                    )
    if dropped:
        prs.save(template_path)
    return dropped


def main() -> int:
    template_path = TEMPLATE
    if len(sys.argv) > 1:
        template_path = Path(sys.argv[1])
    if not template_path.exists():
        print(f"FATAL: template introuvable: {template_path}", file=sys.stderr)
        return 1

    print(f"Patching {template_path}...")
    n = drop_extra_pie3d_series(template_path)
    if n == 0:
        print("  Rien à faire (déjà propre).")
    else:
        print(f"\n{n} série(s) supprimée(s). Template sauvé.")

    # Sanity check post-patch
    print("\nVérification post-patch :")
    prs = Presentation(template_path)
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_chart:
                continue
            sers = sh.chart._chartSpace.findall(".//c:pie3DChart/c:ser", NS)
            ok = "OK" if len(sers) == 1 else f"BAD ({len(sers)})"
            print(f"  slide {i} | {sh.name} | {len(sers)} série | {ok}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
