"""
One-shot script : injecte le bandeau PIC (ligne bleue + logo + 3 rectangles
accents verticaux) directement dans les slides Q1-Q10 + remarques + souhaits
de templates/analysis-block.pptx.

POURQUOI :
    Lors du clone cross-PPTX (`_clone_slide_external` dans inject_slides.py),
    le slideLayout/slideMaster source est SKIPPÉ (whitelist /chart, /image…),
    donc la slide arrive sous le layout "Vide" du PPTX formateur cible : on
    perd la ligne horizontale + le logo PIC qui vivaient sur slideMaster3.

SOLUTION :
    Embarquer ces décorations directement dans le spTree de chaque slide
    concernée (Q1-Q10 + remarques + souhaits), donc elles sont clonées avec
    la slide. C'est exactement le même asset que possède déjà PIC sur leur
    template original (image1.jpeg du master3 = pic-logo.jpeg).

USAGE :
    python scripts/patch_analysis_block_bandeau.py

    Idempotent : si le shape "PIC_Bandeau_Line" existe déjà sur la slide,
    on skip cette slide (re-runs sans dégât).

CIBLE : slides 0-based [3..14] = Q1 (s4) → Souhaits (s15).
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches

try:
    from pptx.oxml import parse_xml
except ImportError:  # older python-pptx fallback
    from pptx.oxml.ns import nsmap as _nsmap
    def parse_xml(s):
        return etree.fromstring(s)

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "templates" / "analysis-block.pptx"
LOGO = ROOT / "templates" / "pic-logo.jpeg"

# Slides cibles (0-based) : Q1=3, Q2=4, …, Q10=12, Remarques=13, Souhaits=14
TARGET_SLIDES = list(range(3, 15))

# Coordonnées EMU extraites du golden (slideMaster3.xml)
# 1 inch = 914400 EMU
LINE = {
    "off_x": 622300,    # 0.68"
    "off_y": 687388,    # 0.75"
    "ext_cx": 10769600, # 11.78"
    "ext_cy": 0,
    "color": "0070C0",  # bleu PIC
    "weight": 19050,    # 1.5pt
}

# Logo PIC (top-right)
LOGO_GEOM = {
    "off_x": 11470314,  # 12.54"
    "off_y": 20203,     # 0.02"
    "ext_cx": 721686,   # 0.79"
    "ext_cy": 720000,   # 0.79"
}

# Rectangles accents verticaux gauche (3 segments empilés)
RECTS = [
    # Top : orange accent1 (~FF6600 d'après le master, mais on prend FF6600 pour cohérence avec Rect 7)
    {"name": "PIC_Bandeau_Rect_Top", "off_x": 0, "off_y": -9128,   "ext_cx": 152400, "ext_cy": 2295128, "color": "FF6600"},
    # Middle : bleu PIC 0070C0
    {"name": "PIC_Bandeau_Rect_Mid", "off_x": 0, "off_y": 2286000, "ext_cx": 152400, "ext_cy": 2286000, "color": "0070C0"},
    # Bottom : orange FF6600
    {"name": "PIC_Bandeau_Rect_Bot", "off_x": 0, "off_y": 4560886, "ext_cx": 152400, "ext_cy": 2286000, "color": "FF6600"},
]

LINE_NAME = "PIC_Bandeau_Line"
LOGO_NAME = "PIC_Bandeau_Logo"


def has_shape_named(slide, name: str) -> bool:
    for shape in slide.shapes:
        if shape.name == name:
            return True
    return False


def make_line_xml(shape_id: int) -> etree._Element:
    """Construit un <p:sp prst='line'> bleu horizontal (Line 8 du master3)."""
    xml = f"""
<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{LINE_NAME}"/>
    <p:cNvSpPr><a:spLocks noChangeShapeType="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr bwMode="auto">
    <a:xfrm>
      <a:off x="{LINE['off_x']}" y="{LINE['off_y']}"/>
      <a:ext cx="{LINE['ext_cx']}" cy="{LINE['ext_cy']}"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="{LINE['weight']}">
      <a:solidFill><a:srgbClr val="{LINE['color']}"/></a:solidFill>
      <a:round/>
      <a:headEnd/>
      <a:tailEnd/>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="fr-FR"/></a:p>
  </p:txBody>
</p:sp>"""
    return parse_xml(xml.strip())


def make_rect_xml(shape_id: int, spec: dict) -> etree._Element:
    xml = f"""
<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{spec['name']}"/>
    <p:cNvSpPr><a:spLocks noChangeArrowheads="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr bwMode="auto">
    <a:xfrm>
      <a:off x="{spec['off_x']}" y="{spec['off_y']}"/>
      <a:ext cx="{spec['ext_cx']}" cy="{spec['ext_cy']}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{spec['color']}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="none" anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="fr-FR"/></a:p>
  </p:txBody>
</p:sp>"""
    return parse_xml(xml.strip())


def next_shape_id(slide) -> int:
    """Génère un id unique pour un nouveau shape dans la slide."""
    used = []
    for el in slide.shapes._spTree.iter():
        sid = el.get("id")
        if sid and sid.isdigit():
            used.append(int(sid))
    return (max(used) + 1) if used else 100


def patch_slide(slide, slide_idx: int):
    """Ajoute le bandeau (ligne, logo, 3 rects) à une slide si pas déjà présent."""
    sp_tree = slide.shapes._spTree
    actions = []

    # 1) Ligne horizontale (en premier dans z-order pour être derrière les contenus)
    if not has_shape_named(slide, LINE_NAME):
        sid = next_shape_id(slide)
        line_el = make_line_xml(sid)
        sp_tree.append(line_el)
        actions.append(f"line(id={sid})")

    # 2) 3 rectangles accents verticaux gauche
    for spec in RECTS:
        if not has_shape_named(slide, spec["name"]):
            sid = next_shape_id(slide)
            rect_el = make_rect_xml(sid, spec)
            sp_tree.append(rect_el)
            actions.append(f"rect:{spec['name']}(id={sid})")

    # 3) Logo PIC en haut à droite (via add_picture pour gérer le rel image)
    if not has_shape_named(slide, LOGO_NAME):
        pic = slide.shapes.add_picture(
            str(LOGO),
            left=LOGO_GEOM["off_x"],
            top=LOGO_GEOM["off_y"],
            width=LOGO_GEOM["ext_cx"],
            height=LOGO_GEOM["ext_cy"],
        )
        pic.name = LOGO_NAME
        actions.append(f"logo(name={LOGO_NAME})")

    print(f"  slide idx={slide_idx} (s{slide_idx+1}) → +{len(actions)} shapes : {', '.join(actions) or '(rien, déjà patché)'}")


def main():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template introuvable : {TEMPLATE}")
    if not LOGO.exists():
        raise SystemExit(f"Logo introuvable : {LOGO}")

    print(f"Patch bandeau PIC sur : {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))
    print(f"Slides cibles : 0-based {TARGET_SLIDES} (= Q1 → Souhaits)")
    for idx in TARGET_SLIDES:
        if idx >= len(prs.slides):
            print(f"  [WARN] slide idx={idx} hors limites ({len(prs.slides)} slides)")
            continue
        patch_slide(prs.slides[idx], idx)

    prs.save(str(TEMPLATE))
    print(f"\nSauvegardé : {TEMPLATE}")


if __name__ == "__main__":
    main()
