"""
Patch branding v3 — fixes bandeau + logo + titre sur les 3 templates Pic.

Bugs corrigés (feedback Pierre 18/05/2026 sur ESAT GARRIC-DTH) :
  1. Bande latérale gauche : `Top` était orange (FF6600) → doit être vert (99CC00).
     Master golden Pic = vert / bleu / orange (top/mid/bot), pas orange/bleu/orange.
  2. Slides 1-3 d'analysis-block (Evaluation / Questions / Globale) sans bandeau
     ni logo : l'ancien `patch_analysis_block_bandeau.py` ciblait seulement les
     slides 3..14 (Q1-Souhaits). On étend à 0..14.
  3. Intro-block (sommaire) sans logo top-right : on ajoute.
  4. Titres en placeholder sans `algn` explicite héritent du master1 qui est
     centré → on force `algn="l"` + Calibri bold sur les titres concernés.
  5. Slide Merci (commercial-closing s0) : positions désordonnées, Merci au-dessus
     de la photo plutôt que sous, pas de logo. On réaligne sur le layout golden
     (photo top / Merci middle / phrase bottom) + ajoute logo.

USAGE :
    python scripts/patch_branding_v3.py

    Idempotent : on remplace les `PIC_Bandeau_*` existants au lieu de skipper
    (sinon les anciennes couleurs FF6600 resteraient en place).
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.slide import Slide

try:
    from pptx.oxml import parse_xml
except ImportError:
    def parse_xml(s):
        return etree.fromstring(s)


NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


ROOT = Path(__file__).resolve().parent.parent
ANALYSIS_TPL = ROOT / "templates" / "analysis-block.pptx"
INTRO_TPL = ROOT / "templates" / "intro-block.pptx"
CLOSING_TPL = ROOT / "templates" / "commercial-closing.pptx"
LOGO_PATH = ROOT / "templates" / "pic-logo.jpeg"


# ---------- Bandeau spec (couleurs golden Pic) ----------

LINE = {
    "off_x": 622300,
    "off_y": 687388,
    "ext_cx": 10769600,
    "ext_cy": 0,
    "color": "0070C0",  # bleu PIC
    "weight": 19050,
}

RECTS = [
    # Top : VERT PIC 99CC00 (= theme accent1 du master3/master5 Pic)
    {"name": "PIC_Bandeau_Rect_Top", "off_x": 0, "off_y": -9128,
     "ext_cx": 152400, "ext_cy": 2295128, "color": "99CC00"},
    # Middle : bleu PIC 0070C0
    {"name": "PIC_Bandeau_Rect_Mid", "off_x": 0, "off_y": 2286000,
     "ext_cx": 152400, "ext_cy": 2286000, "color": "0070C0"},
    # Bottom : orange FF6600
    {"name": "PIC_Bandeau_Rect_Bot", "off_x": 0, "off_y": 4560886,
     "ext_cx": 152400, "ext_cy": 2286000, "color": "FF6600"},
]

LOGO_GEOM = {
    "off_x": 11470314,  # 12.54"
    "off_y": 20203,     # 0.02"
    "ext_cx": 721686,   # 0.79"
    "ext_cy": 720000,   # 0.79"
}

LINE_NAME = "PIC_Bandeau_Line"
LOGO_NAME = "PIC_Bandeau_Logo"

ALL_BANDEAU_NAMES = {LINE_NAME, LOGO_NAME} | {r["name"] for r in RECTS}


# ---------- Helpers ----------

def _find_shape(slide: Slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _remove_shape_named(slide: Slide, name: str) -> bool:
    shape = _find_shape(slide, name)
    if shape is None:
        return False
    shape._element.getparent().remove(shape._element)
    return True


def _next_shape_id(slide: Slide) -> int:
    used = []
    for el in slide.shapes._spTree.iter():
        sid = el.get("id")
        if sid and sid.isdigit():
            used.append(int(sid))
    return (max(used) + 1) if used else 100


def _make_line_xml(shape_id: int) -> etree._Element:
    xml = f"""
<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{LINE_NAME}"/>
    <p:cNvSpPr><a:spLocks noChangeShapeType="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr bwMode="auto">
    <a:xfrm>
      <a:off x="{LINE['off_x']}" y="{LINE['off_y']}"/>
      <a:ext cx="{LINE['ext_cx']}" cy="{LINE['ext_cy']}"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="{LINE['weight']}">
      <a:solidFill><a:srgbClr val="{LINE['color']}"/></a:solidFill>
      <a:round/>
      <a:headEnd/>
      <a:tailEnd/>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="fr-FR"/></a:p>
  </p:txBody>
</p:sp>"""
    return parse_xml(xml.strip())


def _make_rect_xml(shape_id: int, spec: dict) -> etree._Element:
    xml = f"""
<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{spec['name']}"/>
    <p:cNvSpPr><a:spLocks noChangeArrowheads="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr bwMode="auto">
    <a:xfrm>
      <a:off x="{spec['off_x']}" y="{spec['off_y']}"/>
      <a:ext cx="{spec['ext_cx']}" cy="{spec['ext_cy']}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{spec['color']}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="none" anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="fr-FR"/></a:p>
  </p:txBody>
</p:sp>"""
    return parse_xml(xml.strip())


def apply_bandeau(slide: Slide, *, with_logo: bool) -> list[str]:
    """Remplace les PIC_Bandeau_* sur la slide. Idempotent (drop+re-add)."""
    actions = []

    # 1. drop tout PIC_Bandeau_* existant
    for name in (LINE_NAME, *[r["name"] for r in RECTS], LOGO_NAME):
        if _remove_shape_named(slide, name):
            actions.append(f"removed:{name}")

    sp_tree = slide.shapes._spTree

    # 2. Line horizontale (bleue, header)
    sid = _next_shape_id(slide)
    sp_tree.append(_make_line_xml(sid))
    actions.append(f"line(id={sid})")

    # 3. 3 rectangles accents verticaux
    for spec in RECTS:
        sid = _next_shape_id(slide)
        sp_tree.append(_make_rect_xml(sid, spec))
        actions.append(f"rect:{spec['name']}={spec['color']}")

    # 4. Logo top-right (optionnel)
    if with_logo:
        pic = slide.shapes.add_picture(
            str(LOGO_PATH),
            left=LOGO_GEOM["off_x"],
            top=LOGO_GEOM["off_y"],
            width=LOGO_GEOM["ext_cx"],
            height=LOGO_GEOM["ext_cy"],
        )
        pic.name = LOGO_NAME
        actions.append(f"logo")

    return actions


# ---------- Titre placeholder force left-align + Calibri ----------

def force_title_left_align(slide: Slide, *, font_size: int = 2400,
                           shape_name: str | None = None) -> bool:
    """Force algn='l' + Calibri bold sur le titre de la slide.

    Cible :
      - `<p:sp>` avec `<p:ph type="title"/>` (cas par défaut), OU
      - `<p:sp>` avec name=`shape_name` si fourni (cas TEXT_BOX nommé
        explicitement comme 'ZoneTexte 3' utilisé pour certains titres).

    Effet : casse l'héritage centré qu'on récupère du master1 'Vide' après
    cross-PPTX clone + force la police Calibri / taille / bold sur chaque run.
    """
    sp_tree = slide.shapes._spTree
    P_NS = f"{{{NS_P}}}"
    A_NS = f"{{{NS_A}}}"

    title_sp = None
    for sp in sp_tree.findall(f"{P_NS}sp"):
        if shape_name is not None:
            cn = sp.find(f"{P_NS}nvSpPr/{P_NS}cNvPr")
            if cn is not None and cn.get("name") == shape_name:
                title_sp = sp
                break
        else:
            ph = sp.find(f"{P_NS}nvSpPr/{P_NS}nvPr/{P_NS}ph")
            if ph is not None and ph.get("type") == "title":
                title_sp = sp
                break
    if title_sp is None:
        return False

    txBody = title_sp.find(f"{P_NS}txBody")
    if txBody is None:
        return False

    for p in txBody.findall(f"{A_NS}p"):
        pPr = p.find(f"{A_NS}pPr")
        if pPr is None:
            pPr = etree.SubElement(p, f"{A_NS}pPr")
            # pPr doit être en premier dans <a:p>
            p.insert(0, pPr)
        pPr.set("algn", "l")

        # Force font sur chaque run
        for r in p.findall(f"{A_NS}r"):
            rPr = r.find(f"{A_NS}rPr")
            if rPr is None:
                rPr = etree.SubElement(r, f"{A_NS}rPr")
                r.insert(0, rPr)
            rPr.set("sz", str(font_size))
            rPr.set("b", "1")
            rPr.set("lang", rPr.get("lang", "fr-FR"))
            # Force solidFill noir si pas déjà couleur
            if rPr.find(f"{A_NS}solidFill") is None:
                sf = etree.SubElement(rPr, f"{A_NS}solidFill")
                etree.SubElement(sf, f"{A_NS}srgbClr").set("val", "000000")
            # Force latin Calibri
            latin = rPr.find(f"{A_NS}latin")
            if latin is None:
                latin = etree.SubElement(rPr, f"{A_NS}latin")
            latin.set("typeface", "Calibri")

    return True


# ---------- Merci slide reposition (commercial-closing slide 0) ----------

def reposition_merci_slide(slide: Slide) -> list[str]:
    """Reorganise la slide Merci : photo top, Merci middle, phrase bottom.

    Match le layout golden (slide formateur 'doublon' Merci). Positions cibles :
      - Image équipe Pic : off=(3526126, 1127966) ext=(4813300, 4013200)
      - Merci !          : off=(5083373, 3909306) ext=(2668551, 1107996)
      - Phrase           : off=(640188, 5324072)  ext=(10585176, 641971)

    On déplace les shapes existantes (par texte / type) — pas de re-creation.
    """
    actions = []
    P_NS = f"{{{NS_P}}}"
    A_NS = f"{{{NS_A}}}"

    # Needle = sous-chaîne sans apostrophe pour matcher quelle que soit la forme
    # (ASCII ' vs typographique ’).
    targets_by_text = [
        ("PIC FORMATION est intervenue", (640188, 5324072, 10585176, 641971)),
        ("Merci",                       (5083373, 3909306, 2668551, 1107996)),
    ]
    group_target = (3526126, 1127966, 4813300, 4013200)

    for shape in slide.shapes:
        # Compose le texte agrégé du shape (top-level seulement, suffit ici)
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text

        moved = False
        for needle, (x, y, cx, cy) in targets_by_text:
            if needle in text:
                shape.left = x
                shape.top = y
                shape.width = cx
                shape.height = cy
                actions.append(f"moved-text:{needle}")
                moved = True
                break
        if moved:
            continue

        # Groupe qui contient image équipe Pic + Merci (Merci est *à l'intérieur*
        # du groupe, donc on le déplacera avec). On déplace donc le groupe.
        if shape.name.startswith("Groupe"):
            shape.left = group_target[0]
            shape.top = group_target[1]
            shape.width = group_target[2]
            shape.height = group_target[3]
            actions.append(f"moved-group:{shape.name}")
            # Le Merci à l'intérieur du groupe garde sa position relative au
            # groupe (pas besoin de matcher séparément).

    return actions


# ---------- Per-template orchestration ----------

def patch_analysis_block():
    """analysis-block.pptx (15 slides) :
       - slides 3..14 (Q1..Q10 + Remarques + Souhaits) : bandeau + logo.
       - slides 0, 1, 2 (Evaluation / Questions fiche / Évaluation globale) :
         AUCUN bandeau — ce sont des slides titre de section, branding Pic
         déjà présent (colibri central etc.), notre bandeau y serait parasite
         (cf. feedback Inès 20/05).
       - slides 13, 14 : force titre left-align (Remarques / Souhaits).
       - slide Q4 (idx 6) : vire le `ZoneTexte 5` pré-baked (artefact CAMARINES
         hérité) qui forçait l'inject à reposer les commentaires sur une zone
         étroite et décalée vs les autres Q. Sans ce shape, inject_slides crée
         une textbox fraîche à (0.5", 5.4", 12.5"×2.0") cohérente avec Q1..Q10.
    """
    print(f"\n=== Patch {ANALYSIS_TPL.name} ===")
    prs = Presentation(str(ANALYSIS_TPL))

    # Indices 0-based : 0=Evaluation, 1=Questions, 2=Globale, 3..12=Q1..Q10, 13=Remarques, 14=Souhaits
    SLIDES_NO_BANDEAU = {0, 1, 2}
    SLIDES_TITLE_LEFT = {13, 14}

    # Slide Q4 (idx 6) : drop `ZoneTexte 5` pré-baked → inject_slides recréera
    # une textbox commentaires standard comme pour les autres Q.
    SLIDE_Q4 = 6

    for idx, slide in enumerate(prs.slides):
        actions = []
        if idx in SLIDES_NO_BANDEAU:
            # Drop d'éventuels bandeaux laissés par une version antérieure.
            for name in (LINE_NAME, *[r["name"] for r in RECTS], LOGO_NAME):
                if _remove_shape_named(slide, name):
                    actions.append(f"removed:{name}")
        else:
            actions.extend(apply_bandeau(slide, with_logo=True))
        if idx in SLIDES_TITLE_LEFT:
            font_sz = 2400  # 24pt
            if force_title_left_align(slide, font_size=font_sz):
                actions.append(f"title-left+Calibri{font_sz//100}pt")
        if idx == SLIDE_Q4:
            if _remove_shape_named(slide, "ZoneTexte 5"):
                actions.append("dropped-Q4-baked-commentaires-box")
        print(f"  slide {idx}: {', '.join(actions) or '(noop)'}")

    prs.save(str(ANALYSIS_TPL))
    print(f"  → saved")


def patch_intro_block():
    """intro-block.pptx (1 slide sommaire) :
       - bandeau + logo
       - titre 'Sommaire' left-align + Calibri 24pt bold (24pt depuis le
         feedback Inès 20/05 — 28pt cassait la cohérence visuelle avec les
         autres titres de section).
    """
    print(f"\n=== Patch {INTRO_TPL.name} ===")
    prs = Presentation(str(INTRO_TPL))
    slide = prs.slides[0]
    actions = apply_bandeau(slide, with_logo=True)
    if force_title_left_align(slide, font_size=2400):
        actions.append("title-left+Calibri24pt")
    print(f"  slide 0: {', '.join(actions)}")
    prs.save(str(INTRO_TPL))
    print(f"  → saved")


def patch_commercial_closing():
    """commercial-closing.pptx (6 slides) :
       - AUCUN bandeau sur les 6 slides — ce sont des slides marketing Pic
         préexistantes (équipe, certifs, catalogue, contact, 2 murs de logos
         partenaires), branding déjà natif, notre bandeau y est parasite
         (cf. feedback Inès 20/05).
       - slide 0 (Merci) : reposition selon layout golden.
       - Drop d'éventuels bandeaux laissés par une version antérieure.
    """
    print(f"\n=== Patch {CLOSING_TPL.name} ===")
    prs = Presentation(str(CLOSING_TPL))
    bandeau_names = (LINE_NAME, *[r["name"] for r in RECTS], LOGO_NAME)
    for idx, slide in enumerate(prs.slides):
        actions = []
        for name in bandeau_names:
            if _remove_shape_named(slide, name):
                actions.append(f"removed:{name}")
        if idx == 0:
            moves = reposition_merci_slide(slide)
            actions.extend(moves)
        print(f"  slide {idx}: {', '.join(actions) or '(noop)'}")
    prs.save(str(CLOSING_TPL))
    print(f"  → saved")


def main():
    for path in (ANALYSIS_TPL, INTRO_TPL, CLOSING_TPL, LOGO_PATH):
        if not path.exists():
            raise SystemExit(f"Manque : {path}")

    patch_analysis_block()
    patch_intro_block()
    patch_commercial_closing()

    print("\n[OK] Branding v3 appliqué sur les 3 templates.")


if __name__ == "__main__":
    main()
