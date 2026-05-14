"""
One-shot script : construit `templates/intro-block.pptx` (1 slide = sommaire
Pic statique) à partir d'un livrable client final Pic.

POURQUOI :
    Audit 14/05 — la slide 2 (sommaire) est 100 % statique sur les 6 livrables
    Pic analysés (CAMARINES + 5 EXEMPLES) : contenu identique au caractère
    près, aucune mention nom client / titre formation / date. Plutôt que de
    laisser Pierre l'ajouter manuellement dans son template formateur, on
    l'extrait une fois et on l'injecte automatiquement côté Nefia, comme
    on fait déjà pour `commercial-closing.pptx`.

APPROCHE :
    1. Crée un PPTX blank (template par défaut python-pptx)
    2. Clone la slide 2 du livrable source via `_clone_slide_external`
       du module `inject_slides` (whitelist : image + chart cloned, layout
       skip — bandeau ajouté ensuite via patch_intro_block_bandeau.py)
    3. Sauve dans `templates/intro-block.pptx`

USAGE :
    python scripts/build_intro_block_template.py <source-livrable.pptx>

    Exemple :
    python scripts/build_intro_block_template.py \\
      "/Users/guillaume/Downloads/PIC FORMATION PARTAGE NEFIA/2026-03-09-ESAT CAMARINES-SYNTHESE FORMATIVE-GESTIONS DES CONFLITS.pptx"
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
DEST = ROOT / "templates" / "intro-block.pptx"
SLIDE_INDEX_0BASED = 1  # slide 2 = sommaire


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python scripts/build_intro_block_template.py <source.pptx>")
        sys.exit(1)

    source_path = Path(sys.argv[1])
    if not source_path.exists():
        raise FileNotFoundError(f"Source introuvable : {source_path}")

    sys.path.insert(0, str(ROOT))
    from inject_slides import _clone_slide_external

    # Source : livrable Pic complet
    src_prs = Presentation(str(source_path))
    if len(src_prs.slides) <= SLIDE_INDEX_0BASED:
        raise ValueError(
            f"Source ne contient pas de slide à l'index {SLIDE_INDEX_0BASED} "
            f"(n_slides={len(src_prs.slides)})"
        )
    src_slide = src_prs.slides[SLIDE_INDEX_0BASED]

    # Target : PPTX blank
    target_prs = Presentation()
    # Supprimer la slide par défaut si présente (python-pptx Presentation() est vide,
    # mais sécurité au cas où)
    sld_id_lst = target_prs.slides._sldIdLst
    for sid in list(sld_id_lst):
        sld_id_lst.remove(sid)

    # Cloner la slide sommaire
    _clone_slide_external(src_slide, target_prs)

    DEST.parent.mkdir(parents=True, exist_ok=True)
    target_prs.save(str(DEST))

    n = len(target_prs.slides)
    print(f"[OK] intro-block.pptx créé ({n} slide)")
    print(f"     → {DEST}")
    print(f"     Source : {source_path.name}")
    print(f"     Slide  : index {SLIDE_INDEX_0BASED} (= slide {SLIDE_INDEX_0BASED + 1})")
    print()
    print("Prochaine étape : appliquer le bandeau Pic via")
    print("  python scripts/patch_intro_block_bandeau.py")


if __name__ == "__main__":
    main()
