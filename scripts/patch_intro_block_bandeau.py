"""
One-shot script : injecte le bandeau PIC (ligne bleue + 3 rectangles accents
verticaux gauche) dans la slide unique de templates/intro-block.pptx.

POURQUOI :
    Même problème que pour `analysis-block.pptx` (cf. patch_analysis_block_bandeau)
    — le clone cross-PPTX `_clone_slide_external` skippe le slideLayout source,
    donc la slide arrive sous le layout "Vide" du formateur cible et perd la
    ligne horizontale + les 3 rectangles latéraux qui vivaient sur le master.

    NB : le logo PIC est déjà embarqué comme `<p:pic>` direct dans la slide
    sommaire source (Picture 4 du CAMARINES) → préservé nativement au clone,
    pas besoin de l'ajouter.

USAGE :
    python scripts/patch_intro_block_bandeau.py

    Idempotent : si le shape "PIC_Bandeau_Line" existe déjà, on skip.
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation

try:
    from pptx.oxml import parse_xml
except ImportError:
    def parse_xml(s):
        return etree.fromstring(s)

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "templates" / "intro-block.pptx"

# Coordonnées identiques à patch_analysis_block_bandeau.py (même master3 source)
LINE = {
    "off_x": 622300, "off_y": 687388,
    "ext_cx": 10769600, "ext_cy": 0,
    "color": "0070C0", "weight": 19050,
}

RECTS = [
    {"name": "PIC_Bandeau_Rect_Top", "off_x": 0, "off_y": -9128,   "ext_cx": 152400, "ext_cy": 2295128, "color": "FF6600"},
    {"name": "PIC_Bandeau_Rect_Mid", "off_x": 0, "off_y": 2286000, "ext_cx": 152400, "ext_cy": 2286000, "color": "0070C0"},
    {"name": "PIC_Bandeau_Rect_Bot", "off_x": 0, "off_y": 4560886, "ext_cx": 152400, "ext_cy": 2286000, "color": "FF6600"},
]

LINE_NAME = "PIC_Bandeau_Line"


def has_shape_named(slide, name: str) -> bool:
    return any(shape.name == name for shape in slide.shapes)


def make_line_xml(shape_id: int):
    xml = f"""
<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{LINE_NAME}"/>
    <p:cNvSpPr><a:spLocks noChangeShapeType="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr bwMode="auto">
    <a:xfrm>
      <a:off x="{LINE['off_x']}" y="{LINE['off_y']}"/>
      <a:ext cx="{LINE['ext_cx']}" cy="{LINE['ext_cy']}"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="{LINE['weight']}">
      <a:solidFill><a:srgbClr val="{LINE['color']}"/></a:solidFill>
      <a:round/>
      <a:headEnd/>
      <a:tailEnd/>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="fr-FR"/></a:p>
  </p:txBody>
</p:sp>"""
    return parse_xml(xml.strip())


def make_rect_xml(shape_id: int, spec: dict):
    xml = f"""
<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{spec['name']}"/>
    <p:cNvSpPr><a:spLocks noChangeArrowheads="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr bwMode="auto">
    <a:xfrm>
      <a:off x="{spec['off_x']}" y="{spec['off_y']}"/>
      <a:ext cx="{spec['ext_cx']}" cy="{spec['ext_cy']}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{spec['color']}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="none" anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="fr-FR"/></a:p>
  </p:txBody>
</p:sp>"""
    return parse_xml(xml.strip())


def next_shape_id(slide) -> int:
    used = []
    for el in slide.shapes._spTree.iter():
        sid = el.get("id")
        if sid and sid.isdigit():
            used.append(int(sid))
    return (max(used) + 1) if used else 100


def patch_slide(slide):
    sp_tree = slide.shapes._spTree
    actions = []
    if not has_shape_named(slide, LINE_NAME):
        sid = next_shape_id(slide)
        sp_tree.append(make_line_xml(sid))
        actions.append(f"line(id={sid})")
    for spec in RECTS:
        if not has_shape_named(slide, spec["name"]):
            sid = next_shape_id(slide)
            sp_tree.append(make_rect_xml(sid, spec))
            actions.append(f"rect:{spec['name']}(id={sid})")
    print(f"  slide 0 → +{len(actions)} shapes : {', '.join(actions) or '(rien, déjà patché)'}")


def main():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template introuvable : {TEMPLATE} — lancer d'abord build_intro_block_template.py")
    print(f"Patch bandeau PIC sur : {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))
    if not prs.slides:
        raise SystemExit("intro-block.pptx ne contient aucune slide")
    patch_slide(prs.slides[0])
    prs.save(str(TEMPLATE))
    print(f"\nSauvegardé : {TEMPLATE}")


if __name__ == "__main__":
    main()
