"""
Pipeline v3 Pic-Formation : injection de slides d'analyse dans le PPTX formateur.

Architecture :
    formateur.pptx (input client) → injection avant slide "Analyse des résultats" :
        - 15 slides d'analyse (depuis templates/analysis-block.pptx, remplies avec les données OCR)
    → append à la fin :
        - 6 slides closing commercial PIC (depuis templates/commercial-closing.pptx, statiques)

Le formateur garde la maîtrise de son template/identité ; on n'apporte que la valeur
ajoutée IA (analyse questionnaire) + le bloc commercial PIC standardisé.

Usage :
    from inject_slides import inject_pic_analysis
    inject_pic_analysis(
        formateur_pptx_path="…-FORMATEUR.pptx",
        aggregates=…,  # output de extract.aggregate_participants
        analysis_template_path="templates/analysis-block.pptx",
        closing_template_path="templates/commercial-closing.pptx",
        output_path="out.pptx",
    )

Le cœur technique de l'injection est `_clone_slide_external` : copie de slide
entre deux Presentation distinctes (python-pptx ne le fait pas nativement, donc
on copie deep-clone le XML + récursivement les parts liées — chart, image, etc. —
en remappant les rIds dans l'XML cloné).
"""

from __future__ import annotations

import copy
import re
import tempfile
from pathlib import Path
from typing import Any, Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.opc.package import Part
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Pt


# ---------- Constantes ----------

CHART_CATEGORIES = ["Très satisfait", "Satisfait", "Déçu", "Très déçu"]

NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"c": NS_C, "r": NS_R, "p": NS_P, "a": NS_A}

# Layout "Vide" (Blank) — index 6 dans le master Office FR
BLANK_LAYOUT_IDX = 6

ANCHOR_TITLE = "Analyse des résultats"


# ---------- Helpers shape (réutilisés / inspirés de pptx_generator.py) ----------

def _get_shape_by_name(slide: Slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _wrap_quote(commentaire: str) -> str:
    s = commentaire.strip()
    if not s:
        return ""
    if "«" in s or "»" in s:
        return s
    return f"« {s} »"


def _ensure_pt(cache, idx: int, value: str):
    """Force la présence d'un <c:pt idx=N><c:v>value</c:v></c:pt> dans un cache chart."""
    target = None
    for pt in cache.findall("c:pt", NS):
        if int(pt.get("idx")) == idx:
            target = pt
            break
    if target is None:
        target = etree.SubElement(cache, f"{{{NS_C}}}pt")
        target.set("idx", str(idx))
        cache.remove(target)
        inserted = False
        for pt in cache.findall("c:pt", NS):
            if int(pt.get("idx")) > idx:
                pt.addprevious(target)
                inserted = True
                break
        if not inserted:
            cache.append(target)
    v = target.find("c:v", NS)
    if v is None:
        v = etree.SubElement(target, f"{{{NS_C}}}v")
    v.text = value


def _update_pie3d_chart(chart_shape, counts: dict[str, int]):
    """Force les 4 valeurs du chart 3D pie d'une slide question.

    Pour les catégories à 0, on supprime à la fois le `<c:dPt>` (override couleur)
    ET le `<c:dLbl>` (override label) correspondants. Sinon PowerPoint rend des
    artefacts (slice 0 colorée, label "0 (0 %)" planant à côté du camembert) —
    visible quand 100 % des participants ont coché la même case (cas Q10
    "Très satisfait" 8/0/0/0 sur MALECOT).
    """
    chart = chart_shape.chart
    cs = chart._chartSpace
    series = cs.find(".//c:pie3DChart/c:ser", NS)
    if series is None:
        # Certains charts utilisent c:pieChart au lieu de pie3DChart
        series = cs.find(".//c:pieChart/c:ser", NS)
    if series is None:
        return

    num_cache = series.find("c:val/c:numRef/c:numCache", NS)
    str_cache = series.find("c:cat/c:strRef/c:strCache", NS)

    if num_cache is not None:
        ptc = num_cache.find("c:ptCount", NS)
        if ptc is not None:
            ptc.set("val", "4")
        for idx, cat in enumerate(CHART_CATEGORIES):
            _ensure_pt(num_cache, idx, str(counts.get(cat, 0)))

    if str_cache is not None:
        ptc = str_cache.find("c:ptCount", NS)
        if ptc is not None:
            ptc.set("val", "4")
        for idx, cat in enumerate(CHART_CATEGORIES):
            _ensure_pt(str_cache, idx, cat)

    # Supprimer dPt + dLbl pour les catégories à 0 (évite slice fantôme et label "0 (0 %)")
    values_by_idx = {idx: counts.get(cat, 0) for idx, cat in enumerate(CHART_CATEGORIES)}
    total = sum(values_by_idx.values())
    zero_indices = {idx for idx, v in values_by_idx.items() if v == 0}
    for dpt in list(series.findall("c:dPt", NS)):
        idx_el = dpt.find("c:idx", NS)
        if idx_el is not None and int(idx_el.get("val")) in zero_indices:
            series.remove(dpt)
    dlbls = series.find("c:dLbls", NS)
    if dlbls is not None:
        # Désactiver les flags par défaut au niveau série : sinon PowerPoint rend
        # un label DEFAULT "0 ; 0%" (locale-dépendant) pour chaque point val=0
        # qui n'a pas d'override <c:dLbl> explicite. On laisse les overrides per-point
        # piloter entièrement l'affichage via _set_dlbl_text + showVal=0/showPercent=0.
        _set_series_dlbls_flags_off(dlbls)
        # Le manualLayout idx=0 du template est calibré pour le centre VISUEL en
        # cas mono-cat 100% (offset y=-0.45 depuis le centroid du disque entier).
        # En multi-cat, le centroid d'une slice n'est plus au centre du disque
        # → l'offset envoie le label hors slice. Fix : strip le manualLayout +
        # forcer dLblPos=ctr (= centroid de chaque slice) en multi-cat.
        nb_non_zero = sum(1 for v in values_by_idx.values() if v > 0)
        multi_cat = nb_non_zero >= 2
        for dlbl in list(dlbls.findall("c:dLbl", NS)):
            idx_el = dlbl.find("c:idx", NS)
            if idx_el is None:
                continue
            i = int(idx_el.get("val"))
            if i in zero_indices:
                dlbls.remove(dlbl)
            elif total > 0:
                if multi_cat:
                    _strip_dlbl_layout_and_force_ctr(dlbl)
                # Forcer le texte du label : "{value} ({pct} %)" — cohérent avec PR #5
                # (sinon PowerPoint affiche valeur + retour ligne + % à cause de showVal+showPercent).
                _set_dlbl_text(dlbl, f"{values_by_idx[i]} ({round(values_by_idx[i] / total * 100)} %)")


def _strip_dlbl_layout_and_force_ctr(dlbl):
    """En multi-cat : vire le <c:layout> hérité du template + force dLblPos=ctr.

    Le template a un manualLayout idx=0 calibré pour le cas mono-cat 100%
    (offset depuis centroid disque entier). En multi-cat, on veut le label
    au centroid de la slice → dLblPos=ctr sans offset.
    """
    for layout in dlbl.findall("c:layout", NS):
        dlbl.remove(layout)
    existing = dlbl.find("c:dLblPos", NS)
    if existing is not None:
        existing.set("val", "ctr")
        return
    pos = etree.Element(f"{{{NS_C}}}dLblPos")
    pos.set("val", "ctr")
    # Position canonique : avant les c:show* / c:separator / c:extLst
    insert_before = None
    for tag in ("c:showLegendKey", "c:showVal", "c:showCatName", "c:showSerName",
                "c:showPercent", "c:showBubbleSize", "c:separator", "c:extLst"):
        el = dlbl.find(tag, NS)
        if el is not None:
            insert_before = el
            break
    if insert_before is not None:
        insert_before.addprevious(pos)
    else:
        dlbl.append(pos)


def _set_series_dlbls_flags_off(dlbls):
    """Force showVal=0, showPercent=0, showCatName=0 sur le <c:dLbls> SÉRIE.

    Les flags show* au niveau série dictent le rendu DEFAULT des labels pour les
    points sans override <c:dLbl> per-point. Si un point a val=0 sans override,
    PowerPoint affiche "0 ; 0%" (locale-dépendant). En mettant ces flags à 0,
    seuls les points avec un override explicite (avec _set_dlbl_text) afficheront
    un label.

    showLegendKey reste à 0 (cohérent avec template).
    """
    flags_to_zero = {
        "c:showLegendKey": "0",
        "c:showVal": "0",
        "c:showCatName": "0",
        "c:showSerName": "0",
        "c:showPercent": "0",
        "c:showBubbleSize": "0",
    }
    for tag, val in flags_to_zero.items():
        el = dlbls.find(tag, NS)
        if el is not None:
            el.set("val", val)
        else:
            # Crée le flag manquant (rare, mais possible si template incomplet)
            local = tag.split(":")[1]
            new_el = etree.SubElement(dlbls, f"{{{NS_C}}}{local}")
            new_el.set("val", val)


def _set_dlbl_text(dlbl, text: str):
    """Pose un texte explicite sur un <c:dLbl> via <c:tx><c:rich>.

    Place tx à la bonne position canonique (après c:idx + c:layout, avant c:spPr/txPr/show*).
    Force showVal/showPercent à 0 pour éviter que PowerPoint duplique texte + %.
    Supprime aussi un éventuel <c:delete val="1"/> hérité du template (sinon
    le label resterait masqué malgré le texte explicite).
    """
    # Retirer ancien c:tx
    for old_tx in dlbl.findall("c:tx", NS):
        dlbl.remove(old_tx)
    # Retirer un éventuel <c:delete> hérité du template
    for old_del in dlbl.findall("c:delete", NS):
        dlbl.remove(old_del)
    # Construire <c:tx><c:rich><a:bodyPr wrap="none"/><a:lstStyle/>
    #   <a:p><a:pPr><a:buNone/></a:pPr><a:r><a:t>...</a:t></a:r></a:p></c:rich></c:tx>
    # wrap="none" : feedback Inès 13/05 — décocher "Renvoyer le texte à la ligne dans la forme"
    # buNone : sans ça PowerPoint applique le bullet par défaut du chart master (carré
    # blanc/transparent à gauche du label, visible feedback Inès "carré jaune" 13/05)
    tx = etree.Element(f"{{{NS_C}}}tx")
    rich = etree.SubElement(tx, f"{{{NS_C}}}rich")
    body_pr = etree.SubElement(rich, f"{{{NS_A}}}bodyPr")
    body_pr.set("wrap", "none")
    etree.SubElement(rich, f"{{{NS_A}}}lstStyle")
    p = etree.SubElement(rich, f"{{{NS_A}}}p")
    pPr = etree.SubElement(p, f"{{{NS_A}}}pPr")
    etree.SubElement(pPr, f"{{{NS_A}}}buNone")
    r = etree.SubElement(p, f"{{{NS_A}}}r")
    t = etree.SubElement(r, f"{{{NS_A}}}t")
    t.text = text
    # Insertion canonique : juste avant c:spPr, c:txPr, c:dLblPos, ou show* — sinon append
    insert_before = None
    for tag in ("c:spPr", "c:txPr", "c:dLblPos", "c:showLegendKey", "c:showVal",
                "c:showCatName", "c:showSerName", "c:showPercent", "c:showBubbleSize",
                "c:separator", "c:extLst"):
        el = dlbl.find(tag, NS)
        if el is not None:
            insert_before = el
            break
    if insert_before is not None:
        insert_before.addprevious(tx)
    else:
        dlbl.append(tx)
    # Désactiver showVal/showPercent pour ce dLbl (le texte explicite suffit)
    for tag in ("c:showVal", "c:showPercent", "c:showCatName", "c:showSerName"):
        el = dlbl.find(tag, NS)
        if el is not None:
            el.set("val", "0")
    # Forcer spPr noFill + ln noFill si absent — sinon PowerPoint applique le
    # rendu par défaut (fond blanc + bordure noire = "carré blanc" sur slice
    # colorée, feedback Inès 13/05). Le template a noFill explicite seulement
    # sur le dLbl idx=0 ; les autres (qui étaient delete=1) n'ont pas de spPr.
    if dlbl.find("c:spPr", NS) is None:
        spPr = etree.Element(f"{{{NS_C}}}spPr")
        etree.SubElement(spPr, f"{{{NS_A}}}noFill")
        ln = etree.SubElement(spPr, f"{{{NS_A}}}ln")
        etree.SubElement(ln, f"{{{NS_A}}}noFill")
        # Position canonique : juste après c:tx, avant c:txPr/c:dLblPos/c:show*
        insert_before = None
        for tag in ("c:txPr", "c:dLblPos", "c:showLegendKey", "c:showVal",
                    "c:showCatName", "c:showSerName", "c:showPercent",
                    "c:showBubbleSize", "c:separator", "c:extLst"):
            el = dlbl.find(tag, NS)
            if el is not None:
                insert_before = el
                break
        if insert_before is not None:
            insert_before.addprevious(spPr)
        else:
            dlbl.append(spPr)


# Couleur texte body PIC (gris très foncé, presque noir — extraite des slides golden)
_PIC_BODY_COLOR = RGBColor(0x20, 0x21, 0x24)


def _style_run(run, *, size_pt: int, bold: bool = False,
               color: RGBColor = _PIC_BODY_COLOR, font_name: str = "Calibri"):
    """Force explicitement la police, taille, gras et couleur sur un run.

    Indispensable pour les contenus injectés dans des slides clonées cross-PPTX :
    l'inheritance master/layout est perdue lors du clone (cf. memory
    `feedback_pptx_cross_merge`), donc tout doit être posé au niveau run.
    """
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    try:
        run.font.color.rgb = color
    except Exception:
        # Couleur déjà héritée d'un thème scheme — ignore silencieux
        pass


def _apply_pic_bullet(paragraph):
    """Force le style de puce PIC sur un paragraphe : Wingdings 'q' bleu 0070C0.

    Reproduit exactement le pPr du template `analysis-block.pptx` slide 13 (Vos
    remarques), pour que les paragraphes dynamiques injectés via `add_paragraph()`
    soient rendus avec les mêmes puces que les slides statiques. Sans ça, les puces
    héritent du layout, qui est perdu cross-PPTX clone (cf. memory
    `feedback_pptx_cross_merge_fonts`).

    Idempotent : remplace marL/indent/algn s'ils existent, et toute conf bullet
    précédente (buNone, buAutoNum, buChar) est retirée avant insertion.
    """
    pPr = paragraph._pPr
    if pPr is None:
        # Crée le pPr en première position dans <a:p>
        p = paragraph._p
        pPr = etree.SubElement(p, f"{{{NS_A}}}pPr")
        p.remove(pPr)
        p.insert(0, pPr)

    pPr.set("marL", "285750")
    pPr.set("indent", "-285750")
    pPr.set("algn", "l")

    # Purger les éléments bullet existants pour éviter doublons (buNone bloquerait l'affichage)
    for tag in ("buNone", "buAutoNum", "buChar", "buFont", "buClr", "buSzPct", "buSzPts"):
        for el in pPr.findall(f"{{{NS_A}}}{tag}"):
            pPr.remove(el)

    # Reconstruire l'ordre canonique : lnSpc → buClr → buSzPct → buFont → buChar
    lnSpc = pPr.find(f"{{{NS_A}}}lnSpc")
    if lnSpc is None:
        lnSpc = etree.SubElement(pPr, f"{{{NS_A}}}lnSpc")
        pPr.remove(lnSpc)
        pPr.insert(0, lnSpc)
        spcPct = etree.SubElement(lnSpc, f"{{{NS_A}}}spcPct")
        spcPct.set("val", "150000")

    buClr = etree.SubElement(pPr, f"{{{NS_A}}}buClr")
    srgbClr = etree.SubElement(buClr, f"{{{NS_A}}}srgbClr")
    srgbClr.set("val", "0070C0")

    buSzPct = etree.SubElement(pPr, f"{{{NS_A}}}buSzPct")
    buSzPct.set("val", "75000")

    buFont = etree.SubElement(pPr, f"{{{NS_A}}}buFont")
    buFont.set("typeface", "Wingdings")
    buFont.set("panose", "05000000000000000000")
    buFont.set("pitchFamily", "2")
    buFont.set("charset", "2")

    buChar = etree.SubElement(pPr, f"{{{NS_A}}}buChar")
    buChar.set("char", "q")


def _clear_paragraph_bullet(paragraph):
    """Marque explicitement un paragraphe SANS puce (buNone).

    Utile pour les paragraphes 'header' (ex: 'Mes commentaires :') au-dessus
    d'une liste à puces — sinon ils hériteraient de la puce du paragraphe suivant
    ou du layout.
    """
    pPr = paragraph._pPr
    if pPr is None:
        p = paragraph._p
        pPr = etree.SubElement(p, f"{{{NS_A}}}pPr")
        p.remove(pPr)
        p.insert(0, pPr)
    for tag in ("buAutoNum", "buChar"):
        for el in pPr.findall(f"{{{NS_A}}}{tag}"):
            pPr.remove(el)
    if pPr.find(f"{{{NS_A}}}buNone") is None:
        buNone = etree.SubElement(pPr, f"{{{NS_A}}}buNone")


def _set_textbox_with_prenoms(shape, items: list[dict], *, font_size_pt: int = 18):
    """Pose une liste 'Prénom : « commentaire »' dans un textbox, prénom en gras,
    avec puces PIC (Wingdings 'q' bleu) sur chaque item.

    Force Calibri + taille + couleur sur chaque run (pas d'héritage cross-PPTX).
    Default 18pt = taille body lvl1 du master PIC original.
    """
    tf = shape.text_frame
    tf.clear()
    if not items:
        return
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run_prenom = p.add_run()
        run_prenom.text = item["prenom"]
        _style_run(run_prenom, size_pt=font_size_pt, bold=True)
        run_rest = p.add_run()
        run_rest.text = f' : {_wrap_quote(item.get("commentaire", ""))}'
        _style_run(run_rest, size_pt=font_size_pt, bold=False)
        _apply_pic_bullet(p)


def _set_commentaires_textbox(shape, commentaires: list[dict], *, font_size_pt: int = 16):
    """Pose 'Mes commentaires :' (gras, sans puce) puis chaque ligne Prénom : « ... »
    avec puces PIC.

    Force Calibri + taille + couleur sur chaque run. Default 16pt (feedback Inès
    13/05 : minimum 16, plus 14).
    """
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r = p1.add_run()
    r.text = "Mes commentaires :"
    _style_run(r, size_pt=font_size_pt, bold=True)
    _clear_paragraph_bullet(p1)
    for item in commentaires:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = item["prenom"]
        _style_run(r1, size_pt=font_size_pt, bold=True)
        r2 = p.add_run()
        r2.text = f' : {_wrap_quote(item.get("commentaire", ""))}'
        _style_run(r2, size_pt=font_size_pt, bold=False)
        _apply_pic_bullet(p)


# ---------- Anchor lookup ----------

def find_anchor_slide(prs: PresentationType, title: str = ANCHOR_TITLE) -> Optional[int]:
    """Trouve l'index 0-based de la slide dont un shape contient ce texte de titre.

    Match insensible à la casse + normalisation des apostrophes/espaces.
    Renvoie None si pas trouvé.
    """
    needle = _normalize_title(title)
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = _normalize_title(shape.text_frame.text)
            if needle in txt:
                return idx
    return None


def _normalize_title(s: str) -> str:
    return (
        s.lower()
        .replace("’", "'")  # smart apostrophe → straight
        .replace(" ", " ")  # nbsp → space
        .strip()
    )


# ---------- Cross-PPTX slide cloning ----------

# Whitelist des reltypes qu'on clone récursivement (cross-PPTX).
# Tout le reste est ignoré : on ne veut pas cloner les notesSlide
# (back-référence vers la slide → boucle), slideLayout, slideMaster,
# theme — l'inheritance master-level vient de la target_layout.
_CLONABLE_RELTYPES_SUFFIX = {
    "/chart",            # graphique
    "/image",            # image (PNG/JPG/GIF…)
    "/oleObject",        # objet OLE embed
    "/package",          # XLSX embarqué d'un chart
    "/diagramData",      # SmartArt
    "/diagramLayout",
    "/diagramQuickStyle",
    "/diagramColors",
    "/audio",
    "/video",
    "/media",
}


def _is_clonable(reltype: str) -> bool:
    return any(reltype.endswith(suffix) for suffix in _CLONABLE_RELTYPES_SUFFIX)


def _clone_part(src_part, target_pkg, cache: dict):
    """Clone récursif (whitelist) d'un Part vers un autre package OPC.

    `cache` (id(src_part) → new_part) déduplique les parts partagés
    (ex: même logo référencé par 5 slides → 1 seul clone).
    Retourne le nouveau Part.
    """
    if id(src_part) in cache:
        return cache[id(src_part)]

    template = _partname_template(src_part)
    partname = target_pkg.next_partname(template)
    new_part = Part(partname, src_part.content_type, target_pkg, blob=src_part.blob)
    cache[id(src_part)] = new_part

    for rel in src_part.rels.values():
        if rel.is_external:
            new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        elif _is_clonable(rel.reltype):
            new_sub = _clone_part(rel.target_part, target_pkg, cache)
            new_part.relate_to(new_sub, rel.reltype)
        # Sinon : skip silencieux (slideLayout/Master/theme/notes/tags…)

    return new_part


def _partname_template(part) -> str:
    """De '/ppt/charts/chart7.xml' produit '/ppt/charts/chart%d.xml'."""
    name = str(part.partname)
    # Remplace la 1re séquence de chiffres par %d (en partant de la fin)
    m = re.search(r"\d+(?=\.[^/]+$)", name)
    if m:
        return name[: m.start()] + "%d" + name[m.end() :]
    # Fallback : insérer %d avant l'extension
    base, _, ext = name.rpartition(".")
    return f"{base}%d.{ext}"


def _remap_rels_in_element(el, rid_map: dict[str, str]):
    """Remplace tous les attributs r:id, r:embed, r:link selon rid_map."""
    attrs = (f"{{{NS_R}}}id", f"{{{NS_R}}}embed", f"{{{NS_R}}}link")
    for sub in el.iter():
        for attr in attrs:
            v = sub.get(attr)
            if v is not None and v in rid_map:
                sub.set(attr, rid_map[v])


def _clone_slide_external(src_slide: Slide, target_prs: PresentationType,
                          target_layout=None, *, part_cache: dict | None = None) -> Slide:
    """Clone une slide depuis une autre Presentation vers target_prs.

    Approche :
      1. Crée une nouvelle slide blanche dans target_prs (layout `Vide`).
      2. Vide les placeholders du layout sur la nouvelle slide.
      3. Pour chaque relationship de la slide source, clone le part associé
         dans le package cible et garde un mapping rId_src → rId_dest.
      4. Deep-copy chaque shape XML du src, remappe les rIds, append au spTree.

    NB : on perd l'inheritance du slideLayout source (couleurs/polices du master).
    Pour notre cas (template PIC stable, formatage shape-level), c'est OK.
    """
    if target_layout is None:
        target_layout = target_prs.slide_layouts[BLANK_LAYOUT_IDX]

    new_slide = target_prs.slides.add_slide(target_layout)

    # Vider les placeholders (le layout vide n'en a normalement pas, mais sécurité)
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)

    src_part = src_slide.part
    new_part = new_slide.part
    target_pkg = target_prs.part.package
    if part_cache is None:
        part_cache = {}

    # Map rIds source → cibles (whitelist : on ne clone que chart/image/etc.)
    rid_map: dict[str, str] = {}
    for src_rId, src_rel in src_part.rels.items():
        if src_rel.is_external:
            new_rId = new_part.relate_to(src_rel.target_ref, src_rel.reltype, is_external=True)
            rid_map[src_rId] = new_rId
        elif _is_clonable(src_rel.reltype):
            new_sub_part = _clone_part(src_rel.target_part, target_pkg, part_cache)
            new_rId = new_part.relate_to(new_sub_part, src_rel.reltype)
            rid_map[src_rId] = new_rId
        # Else: slideLayout, notesSlide, tags… → skip
        # Le shape XML pourrait référencer un rId skippé ; on laissera tel quel.
        # En pratique, les shapes ne référencent que les rels whitelisted (charts, images).

    # Copier les shapes (deep-copy puis remap rIds)
    sp_tree = new_slide.shapes._spTree
    src_sp_tree = src_slide.shapes._spTree
    # Le spTree du src contient nvGrpSpPr + grpSpPr (header) puis les shapes
    # On copie uniquement les éléments de shape, pas le header (déjà dans new_slide)
    SHAPE_TAGS = {f"{{{NS_P}}}{tag}" for tag in
                  ("sp", "grpSp", "graphicFrame", "cxnSp", "pic", "contentPart")}
    for child in src_sp_tree:
        if child.tag in SHAPE_TAGS:
            new_el = copy.deepcopy(child)
            _remap_rels_in_element(new_el, rid_map)
            sp_tree.append(new_el)

    return new_slide


def inject_external_slides(target_prs: PresentationType, source_pptx_path: str | Path,
                           before_idx: Optional[int] = None) -> int:
    """Clone toutes les slides de `source_pptx_path` dans `target_prs`.

    Si `before_idx` est None, append à la fin.
    Sinon, insère avant l'index 0-based donné (les slides après sont décalées).

    Retourne le nombre de slides ajoutées.
    """
    src_prs = Presentation(str(source_pptx_path))
    n_src = len(src_prs.slides)

    # Cache cross-slide pour dédupliquer les parts partagés (logos, etc.)
    part_cache: dict = {}

    # Cloner toutes les slides en fin de target_prs
    for src_slide in src_prs.slides:
        _clone_slide_external(src_slide, target_prs, part_cache=part_cache)

    if before_idx is not None:
        # Réorganiser : déplacer les n_src dernières slides à before_idx
        sld_id_lst = target_prs.slides._sldIdLst
        sld_ids = list(sld_id_lst)
        new_ids = sld_ids[-n_src:]
        for sid in new_ids:
            sld_id_lst.remove(sid)
        # Insérer à before_idx
        for offset, sid in enumerate(new_ids):
            sld_id_lst.insert(before_idx + offset, sid)

    return n_src


# ---------- Slide deletion (low-level, depuis pptx_generator.py) ----------

def _delete_slide(prs: PresentationType, slide_idx_0based: int):
    """Supprime une slide par son index 0-based."""
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    if slide_idx_0based >= len(sld_ids):
        return
    sld_id = sld_ids[slide_idx_0based]
    rId = sld_id.get(f"{{{NS_R}}}id")
    sld_id_lst.remove(sld_id)
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


# ---------- Filling analysis-block.pptx with aggregated data ----------

# Mapping 0-based des slides analysis-block.pptx
# (15 slides : intro/questions/global + Q1-Q10 + remarques/souhaits)
SLIDE_INTRO = 0          # "Evaluation de la formation"
SLIDE_QUESTIONS_LIST = 1 # "Les questions de la fiche..."
SLIDE_GLOBAL = 2         # "Évaluation globale de l'action"
SLIDE_Q_FIRST = 3        # Q1
SLIDE_Q_LAST = 12        # Q10
SLIDE_REMARQUES = 13
SLIDE_SOUHAITS = 14


def _fill_question_slide(slide: Slide, counts: dict[str, int],
                         commentaires: list[dict]):
    """Met à jour le chart pie3D + zone commentaires d'une slide Q1-Q10."""
    # Chart : trouver le shape avec has_chart
    for shape in slide.shapes:
        if shape.has_chart:
            _update_pie3d_chart(shape, counts)
            break

    # Commentaires : trouver/créer le textbox "ZoneTexte 5" (en bas) ou similaire
    # NB: dans analysis-block.pptx, seules certaines slides ont ce textbox
    # (héritage de la slide originale d'EXEMPLE 1). Stratégie :
    # - Chercher un textbox dont le texte commence par "Mes commentaires" ou nom "ZoneTexte 5".
    # - Si absent et commentaires vides : skip.
    # - Si absent et commentaires non-vides : créer un textbox.
    target = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip().lower()
        if txt.startswith("mes commentaires") or shape.name == "ZoneTexte 5":
            target = shape
            break

    if commentaires:
        if target is None:
            # Créer un textbox de commentaires en bas de la slide
            from pptx.util import Inches
            target = slide.shapes.add_textbox(Inches(0.5), Inches(5.4),
                                              Inches(12.5), Inches(2.0))
        _set_commentaires_textbox(target, commentaires)
    else:
        # Pas de commentaires : vider le textbox existant si présent
        if target is not None:
            target.text_frame.clear()


def _fill_libre_slide(slide: Slide, items: list[dict]):
    """S14/S15 : pose 'Prénom : « commentaire »' dans le textbox principal."""
    # Le textbox principal est généralement nommé "ZoneTexte 1"
    target = _get_shape_by_name(slide, "ZoneTexte 1")
    if target is None:
        # Fallback : 1er textbox non-titre, non-numéro de slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name.startswith("Espace réservé") or shape.name == "Rectangle 2":
                continue
            target = shape
            break
    if target is None:
        return
    if not items:
        target.text_frame.clear()
        return
    _set_textbox_with_prenoms(target, items)


def fill_analysis_block(prs: PresentationType, aggregates: dict[str, Any]):
    """Remplit les 15 slides d'analysis-block.pptx avec les données agrégées.

    `aggregates` est l'output de `extract.aggregate_participants` :
        {
            "nb_participants": int,
            "questions": list[10] of {counts, commentaires},
            "remarques_libres": list[{prenom, commentaire}],
            "souhaits_libres": list[{prenom, commentaire}],
        }
    """
    # S1, S2, S3 = intros statiques, on ne touche pas pour l'instant
    # (texte générique du template, pas de variable formateur)

    # Q1-Q10
    questions = aggregates.get("questions", [])
    for i, q_data in enumerate(questions[:10]):
        slide_idx = SLIDE_Q_FIRST + i
        if slide_idx >= len(prs.slides):
            break
        _fill_question_slide(
            prs.slides[slide_idx],
            counts=q_data.get("counts", {}),
            commentaires=q_data.get("commentaires", []),
        )

    # Remarques libres
    if SLIDE_REMARQUES < len(prs.slides):
        _fill_libre_slide(prs.slides[SLIDE_REMARQUES],
                          aggregates.get("remarques_libres", []))

    # Souhaits libres
    if SLIDE_SOUHAITS < len(prs.slides):
        _fill_libre_slide(prs.slides[SLIDE_SOUHAITS],
                          aggregates.get("souhaits_libres", []))


# ---------- Entry point principal ----------

def inject_pic_analysis(
    formateur_pptx_path: str | Path,
    aggregates: dict[str, Any],
    output_path: str | Path,
    *,
    analysis_template_path: str | Path = "templates/analysis-block.pptx",
    closing_template_path: str | Path = "templates/commercial-closing.pptx",
    anchor_title: str = ANCHOR_TITLE,
) -> Path:
    """Injecte le bloc analyse + closing commercial dans un PPTX formateur PIC.

    Étapes :
      1. Ouvre le PPTX formateur (cible).
      2. Trouve la slide d'ancre "Analyse des résultats" (fallback : position 3).
      3. Ouvre analysis-block.pptx, remplit avec les agrégats, sauve dans un tmpfile.
      4. Clone les 15 slides analyse vers le formateur, AVANT l'ancre.
      5. Clone les 6 slides closing commercial à la fin.
      6. Sauve dans output_path.
    """
    formateur_pptx_path = Path(formateur_pptx_path)
    output_path = Path(output_path)
    analysis_template_path = Path(analysis_template_path)
    closing_template_path = Path(closing_template_path)

    if not formateur_pptx_path.exists():
        raise FileNotFoundError(f"PPTX formateur introuvable : {formateur_pptx_path}")
    if not analysis_template_path.exists():
        raise FileNotFoundError(f"Template analyse introuvable : {analysis_template_path}")
    if not closing_template_path.exists():
        raise FileNotFoundError(f"Template closing introuvable : {closing_template_path}")

    # 1. Ouvrir le formateur
    prs = Presentation(str(formateur_pptx_path))

    # 2. Trouver l'ancre
    anchor_idx = find_anchor_slide(prs, anchor_title)
    if anchor_idx is None:
        # Fallback position 3 (0-based : index 2, après cover + rappels)
        anchor_idx = min(2, len(prs.slides))
        print(f"  [WARN] Ancre '{anchor_title}' non trouvée → fallback position {anchor_idx}")
    else:
        print(f"  [OK] Ancre '{anchor_title}' trouvée à l'index {anchor_idx}")

    # 3. Préparer analysis-block rempli
    analysis_prs = Presentation(str(analysis_template_path))
    fill_analysis_block(analysis_prs, aggregates)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        analysis_filled_path = Path(tmp.name)
    analysis_prs.save(str(analysis_filled_path))

    try:
        # 4. Injecter avant l'ancre
        n_inj = inject_external_slides(prs, analysis_filled_path, before_idx=anchor_idx)
        print(f"  [OK] {n_inj} slides analyse injectées avant l'ancre")

        # 5. Append closing commercial
        n_close = inject_external_slides(prs, closing_template_path, before_idx=None)
        print(f"  [OK] {n_close} slides closing commercial appendées")
    finally:
        analysis_filled_path.unlink(missing_ok=True)

    # 6. Sauver
    prs.save(str(output_path))
    print(f"  [OK] PPTX final → {output_path} ({len(prs.slides)} slides)")
    return output_path


# ---------- CLI ----------

if __name__ == "__main__":
    import json
    import sys

    if len(sys.argv) < 4:
        print("Usage: python inject_slides.py <formateur.pptx> <aggregates.json> <output.pptx>")
        print("       (aggregates.json = output de extract.aggregate_participants)")
        sys.exit(1)

    formateur = sys.argv[1]
    aggregates_path = sys.argv[2]
    output = sys.argv[3]

    with open(aggregates_path) as f:
        aggregates = json.load(f)

    inject_pic_analysis(formateur, aggregates, output)
