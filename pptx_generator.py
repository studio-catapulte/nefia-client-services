"""
Génération PPTX agrégé : N questionnaires → 1 slide par question (camembert + commentaires) + synthèse.
Niveaux de réponse dynamiques (lus depuis les données extraites).
"""

import io
from collections import Counter, defaultdict

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x1F, 0x29, 0x37)
COLOR_LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)

# Palette pour les niveaux de réponse (du plus positif au plus négatif)
LEVEL_COLORS = [
    RGBColor(0x22, 0xC5, 0x5E),  # vert
    RGBColor(0x84, 0xCC, 0x16),  # vert-jaune
    RGBColor(0xFA, 0xCC, 0x15),  # jaune
    RGBColor(0xF9, 0x73, 0x16),  # orange
    RGBColor(0xEF, 0x44, 0x44),  # rouge
]


def _get_color_map(levels: list[str]) -> dict[str, RGBColor]:
    """Associe chaque niveau de réponse à une couleur."""
    colors = {}
    for i, level in enumerate(levels):
        colors[level] = LEVEL_COLORS[min(i, len(LEVEL_COLORS) - 1)]
    return colors


_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def _convert_pie_to_3d(chart) -> None:
    """Convertit un c:pieChart en c:pie3DChart au niveau XML.

    python-pptx 1.x ne sait pas écrire XL_CHART_TYPE.THREE_D_PIE (NotImplementedError),
    mais accepte de créer un PIE classique puis de muter l'élément. PowerPoint lit
    correctement le pie3DChart résultant. Les enfants (series, points, dLbls...) sont
    identiques entre les deux, donc le swap est sûr.
    """
    chartSpace = chart._chartSpace
    pie = chartSpace.find(f".//{{{_C_NS}}}pieChart")
    if pie is not None:
        pie.tag = f"{{{_C_NS}}}pie3DChart"


def _aggregate(all_questionnaires: list[dict]) -> dict:
    """Agrège N questionnaires par question."""
    # Découvrir les niveaux de réponse et les questions
    response_levels = []
    items_order = []
    items_labels = {}
    counts = defaultdict(Counter)  # {question_number: Counter(response -> count)}
    comments = defaultdict(list)   # {question_number: [(participant, comment)]}
    remarks = []
    wishes = []
    metadata = None

    for q in all_questionnaires:
        if not response_levels and q.get("response_levels"):
            response_levels = q["response_levels"]
        if not metadata and q.get("metadata"):
            metadata = q["metadata"]

        participant = q.get("metadata", {}).get("participant", "?")

        for item in q.get("items", []):
            num = item.get("number", 0)
            if num not in items_order:
                items_order.append(num)
                items_labels[num] = item["label"]

            resp = item.get("response", "Non renseigné")
            if resp != "Non renseigné":
                counts[num][resp] += 1

            comment = item.get("comment")
            if comment:
                comments[num].append((participant, comment))

        if q.get("remarks"):
            remarks.append((participant, q["remarks"]))
        if q.get("wishes"):
            wishes.append((participant, q["wishes"]))

    return {
        "response_levels": response_levels or ["Très satisfait", "Satisfait", "Déçu", "Très déçu"],
        "items_order": items_order,
        "items_labels": items_labels,
        "counts": counts,
        "comments": comments,
        "remarks": remarks,
        "wishes": wishes,
        "metadata": metadata or {},
        "n": len(all_questionnaires),
    }


def _add_title_slide(prs, title, subtitle, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_PRIMARY

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"Synthèse de formation"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER

    if subtitle:
        p3 = tf.add_paragraph()
        p3.text = subtitle
        p3.font.size = Pt(16)
        p3.font.color.rgb = COLOR_WHITE
        p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = f"\n{n} participant(s)"
    p4.font.size = Pt(16)
    p4.font.color.rgb = COLOR_WHITE
    p4.alignment = PP_ALIGN.CENTER


def _add_info_slide(prs, metadata, n):
    """Slide rappels de l'intervention."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Rappels de l'intervention"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
    tf2 = info_box.text_frame
    tf2.word_wrap = True

    fields = [
        ("Nombre de participants", str(n)),
        ("Formation", metadata.get("formation", "-")),
        ("Lieu", metadata.get("lieu", "-")),
        ("Formateur", metadata.get("formateur", "-")),
        ("Date", metadata.get("date", "-")),
    ]

    for label, value in fields:
        p = tf2.add_paragraph()
        p.text = f"  {label} : {value}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(6)


def _add_question_slide(prs, number, label, counter, item_comments, color_map, response_levels):
    """1 slide par question : camembert + commentaires."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titre question
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{number}. {label}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # Camembert
    categories = []
    values = []
    for level in response_levels:
        count = counter.get(level, 0)
        if count > 0:
            categories.append(level)
            values.append(count)

    if values:
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("", values)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(0.3), Inches(1.2), Inches(5), Inches(4.5),
            chart_data,
        )
        chart = chart_frame.chart

        series = chart.series[0]
        for idx, cat in enumerate(categories):
            point = series.points[idx]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color_map.get(cat, LEVEL_COLORS[-1])

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)

        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = False
        data_labels.show_value = False
        data_labels.show_category_name = False
        data_labels.font.size = Pt(10)

        total_values = sum(values)
        for idx, value in enumerate(values):
            pct = (value / total_values * 100) if total_values else 0
            point = series.points[idx]
            point.data_label.has_text_frame = True
            point.data_label.text_frame.text = f"{value} ({pct:.0f} %)"

        _convert_pie_to_3d(chart)

    # Commentaires à droite
    if item_comments:
        comment_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4.2), Inches(5))
        tf2 = comment_box.text_frame
        tf2.word_wrap = True

        p_header = tf2.paragraphs[0]
        p_header.text = "Mes commentaires :"
        p_header.font.size = Pt(12)
        p_header.font.bold = True
        p_header.font.color.rgb = COLOR_DARK
        p_header.space_after = Pt(8)

        for participant, comment in item_comments:
            p_c = tf2.add_paragraph()
            p_c.text = f"{participant} : \u00ab {comment} \u00bb"
            p_c.font.size = Pt(10)
            p_c.font.color.rgb = COLOR_DARK
            p_c.space_after = Pt(4)


def _add_global_summary_slide(prs, agg, color_map):
    """Slide synthèse globale avec camembert agrégé."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Evaluation globale de l'action"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    response_levels = agg["response_levels"]
    n_levels = len(response_levels)

    # Agréger toutes les réponses
    global_counter = Counter()
    total = 0
    score_sum = 0

    for num in agg["items_order"]:
        for resp, count in agg["counts"][num].items():
            global_counter[resp] += count
            total += count
            idx = response_levels.index(resp) if resp in response_levels else n_levels - 1
            score = n_levels - idx
            score_sum += score * count

    score_avg = score_sum / total if total else 0
    score_max = n_levels

    # Score
    ratio = score_avg / score_max if score_max else 0
    color = LEVEL_COLORS[0] if ratio >= 0.75 else (LEVEL_COLORS[2] if ratio >= 0.5 else LEVEL_COLORS[-1])

    txBox2 = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(4), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"Score moyen : {score_avg:.1f} / {score_max}"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = color

    p3 = tf2.add_paragraph()
    p3.text = f"{agg['n']} participant(s), {total} reponses"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_DARK

    # Camembert global
    categories = [r for r in response_levels if global_counter.get(r, 0) > 0]
    vals = [global_counter[r] for r in categories]

    if vals:
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("", vals)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(4.5), Inches(1.0), Inches(5), Inches(5),
            chart_data,
        )
        chart = chart_frame.chart

        series = chart.series[0]
        for idx, cat in enumerate(categories):
            point = series.points[idx]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color_map.get(cat, LEVEL_COLORS[-1])

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)

        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_percentage = False
        dl.show_value = False
        dl.show_category_name = False
        dl.font.size = Pt(10)

        total_vals = sum(vals)
        for idx, value in enumerate(vals):
            pct = (value / total_vals * 100) if total_vals else 0
            point = series.points[idx]
            point.data_label.has_text_frame = True
            point.data_label.text_frame.text = f"{value} ({pct:.0f} %)"

        _convert_pie_to_3d(chart)


def _add_text_slide(prs, title, entries):
    """Slide avec texte libre (remarques, souhaits)."""
    if not entries:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True

    for participant, text in entries:
        p_c = tf2.add_paragraph()
        p_c.text = f"{participant} : \u00ab {text} \u00bb"
        p_c.font.size = Pt(11)
        p_c.font.color.rgb = COLOR_DARK
        p_c.space_after = Pt(6)


def generate_pptx(all_questionnaires: list[dict], title: str = "Resultats de satisfaction") -> bytes:
    """Genere un PPTX agrege et retourne les bytes."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    agg = _aggregate(all_questionnaires)
    color_map = _get_color_map(agg["response_levels"])

    meta = agg["metadata"]
    formation = meta.get("formation", title)
    subtitle = f"{meta.get('lieu', '')} - {meta.get('date', '')}".strip(" -")

    _add_title_slide(prs, formation, subtitle, agg["n"])
    _add_info_slide(prs, meta, agg["n"])
    _add_global_summary_slide(prs, agg, color_map)

    # 1 slide par question
    for num in agg["items_order"]:
        label = agg["items_labels"][num]
        counter = agg["counts"][num]
        item_comments = agg["comments"].get(num, [])
        _add_question_slide(prs, num, label, counter, item_comments, color_map, agg["response_levels"])

    _add_text_slide(prs, "Remarques et suggestions", agg["remarks"])
    _add_text_slide(prs, "Souhaits pour d'autres formations", agg["wishes"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
